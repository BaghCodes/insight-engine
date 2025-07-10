import os
import re
import fitz
import docx
import pptx
import openpyxl

UPLOAD_FOLDER = "../../data/uploads"
EXTENSIONS = {".pdf",".xlsx",".docx",".pptx"}

def parse_documents(file_path:str) -> str:
    # detects file type and dispatches to appropriate parser
    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".pdf":
        text_blocks = _parse_pdf(file_path)
    elif extension == ".docx":
        text_blocks = _parse_docx(file_path)
    elif extension == ".pptx":
        text_blocks = _parse_pptx(file_path)
    elif extension == ".xlsx":
        text_blocks = _parse_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {extension}")

    raw_text = "".join(text_blocks)
    return _clean_text(raw_text)

def _parse_pdf(path: str) -> list[str]:
    text_blocks = []
    try:
        doc = fitz.open(path)
        for page in doc:
            text = page.get_text()
            if text:
                text_blocks.append(text)
        doc.close()
    except Exception as e:
        raise RuntimeError(f"Failed to parse PDF: {e}")
    return text_blocks

def _parse_docx(path: str) -> list[str]:
    text_blocks = []
    try:
        document = docx.Document(path)
        for para in document.paragraphs:
            if para.text.strip():
                text_blocks.append(para.text)
    except Exception as e:
        raise RuntimeError(f"Failed to parse DOCX: {e}")
    return text_blocks

def _parse_pptx(path: str) -> list[str]:
    text_blocks = []
    try:
        presentation = pptx.Presentation(path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_blocks.append(shape.text)
    except Exception as e:
        raise RuntimeError(f"Failed to parse PPTX: {e}")
    return text_blocks

def _parse_xlsx(path: str) -> list[str]:
    text_blocks = []
    try:
        workbook = openpyxl.load_workbook(path, data_only=True)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text_blocks.append(str(cell.value))
    except Exception as e:
        raise RuntimeError(f"Failed to parse XLSX: {e}")
    return text_blocks

def _clean_text(text: str) -> str:
    cleaned = re.sub(r"\s+", " ", text)
    cleaned = re.sub(r"[^\x20-\x7E]+", "", cleaned)
    return cleaned.strip()

def parse_uploads():
    # parses all documents in "insight-engine/data/uploads/"
    if not os.path.exists(UPLOAD_FOLDER):
        print(f"Folder not found: {UPLOAD_FOLDER}")
        return

    files_found = False

    for filename in os.listfir(UPLOAD_FOLDER):
        file_path = os.path.join(UPLOAD_FOLDER,filename)

        if not os.path.isfile(file_path) or filename.startswith("."):
            continue
        
        _,ext = os.path.splitext(filename)
        ext = ext.lower()

        if ext in EXTENSIONS:
            files_found=True
            try:
                print(f"Processing {filename}")
                text = parse_document(file_path)
                parsed_output[filename] = text
                print("Text extracted successfully")
            except Exception as e:
                print(f"Error while parsing {filename}")
        else:
            print(f"Skipped unsupported file: {filename}")

    if not files_found:
        print(f"No supported documents found.")
    return parsed_output
if __name__ == "__main__":
    parse_uploads()

